import os
import subprocess
import shutil
from pathlib import Path
import tempfile

from pdf2docx import Converter
from PIL import Image
from PyPDF2 import PdfReader
from reportlab.pdfgen import canvas
import fitz  # PyMuPDF

# Optional imports (SAFE for cloud)
try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except:
    PPTX_AVAILABLE = False


# =====================================================
# 🔧 LIBREOFFICE DETECTION
# =====================================================
def get_soffice_path():
    soffice = shutil.which("soffice")

    if soffice:
        return soffice

    if os.path.exists("/usr/bin/soffice"):  # Streamlit Cloud
        return "/usr/bin/soffice"

    # Windows fallback
    win_paths = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"
    ]
    for p in win_paths:
        if os.path.exists(p):
            return p

    return None


# =====================================================
# 📄 DOCX → PDF
# =====================================================
def docx_to_pdf(input_path, output_path):
    soffice = get_soffice_path()

    if not soffice:
        raise Exception("❌ DOCX → PDF not supported (LibreOffice missing on server)")

    try:
        subprocess.run([
            soffice,
            "--headless",
            "--convert-to", "pdf",
            input_path,
            "--outdir", os.path.dirname(output_path)
        ], check=True)

        generated = Path(input_path).with_suffix(".pdf")

        if not generated.exists():
            raise Exception("Conversion failed")

        os.replace(generated, output_path)
        return output_path

    except Exception as e:
        raise Exception(f"DOCX → PDF failed: {str(e)}")


# =====================================================
# 📊 PPTX → PDF
# =====================================================
def pptx_to_pdf(input_path, output_path):
    soffice = get_soffice_path()

    if not soffice:
        raise Exception("❌ PPTX → PDF not supported (LibreOffice missing)")

    try:
        subprocess.run([
            soffice,
            "--headless",
            "--convert-to", "pdf",
            input_path,
            "--outdir", os.path.dirname(output_path)
        ], check=True)

        generated = Path(input_path).with_suffix(".pdf")

        if not generated.exists():
            raise Exception("Conversion failed")

        os.replace(generated, output_path)
        return output_path

    except Exception as e:
        raise Exception(f"PPTX → PDF failed: {str(e)}")


# =====================================================
# 📄 PDF → DOCX
# =====================================================
def pdf_to_docx(input_path, output_path):
    cv = Converter(input_path)
    cv.convert(output_path)
    cv.close()
    return output_path


# =====================================================
# 📊 PDF → PPTX
# =====================================================
def pdf_to_pptx(input_path, output_path):
    if not PPTX_AVAILABLE:
        raise Exception("❌ python-pptx not installed")

    pdf = fitz.open(input_path)
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    for page in pdf:
        text = page.get_text("text")[:1200]

        slide = prs.slides.add_slide(blank_layout)
        box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(6)
        )
        box.text_frame.text = text

    prs.save(output_path)
    return output_path


# =====================================================
# 🖼 IMAGE → PDF
# =====================================================
def image_to_pdf(input_path, output_path):
    img = Image.open(input_path).convert("RGB")
    img.save(output_path)
    return output_path


# =====================================================
# 🖼 PDF → IMAGE
# =====================================================
def pdf_to_image(input_path, output_path):
    doc = fitz.open(input_path)
    paths = []

    for i, page in enumerate(doc):
        pix = page.get_pixmap()
        out = output_path.replace(".png", f"_{i}.png")
        pix.save(out)
        paths.append(out)

    return paths


# =====================================================
# 📝 TXT → PDF
# =====================================================
def txt_to_pdf(input_path, output_path):
    c = canvas.Canvas(output_path)
    y = 800

    with open(input_path, encoding="utf-8") as f:
        for line in f:
            c.drawString(40, y, line.strip())
            y -= 20

            if y < 50:
                c.showPage()
                y = 800

    c.save()
    return output_path


# =====================================================
# 📄 PDF → TXT
# =====================================================
def pdf_to_txt(input_path, output_path):
    reader = PdfReader(input_path)
    text = "".join([page.extract_text() or "" for page in reader.pages])

    with open(output_path, "w", encoding="utf-8") as f:
        f.write(text)

    return output_path


# =====================================================
# 🔗 MERGE PDFs
# =====================================================
def merge_pdfs(pdf_paths, output_path):
    master = fitz.open()

    for pdf in pdf_paths:
        doc = fitz.open(pdf)
        master.insert_pdf(doc)
        doc.close()

    master.save(output_path)
    return output_path


# =====================================================
# 🔄 CONVERT NON-PDF → PDF
# =====================================================
def to_pdf_if_needed(input_path, temp_dir):
    ext = Path(input_path).suffix.lower()

    if ext == ".pdf":
        return input_path

    pdf_path = os.path.join(temp_dir, Path(input_path).stem + ".pdf")

    conv_map = {
        ".docx": "docx_to_pdf",
        ".pptx": "pptx_to_pdf",
        ".txt": "txt_to_pdf",
        ".jpg": "image_to_pdf",
        ".jpeg": "image_to_pdf",
        ".png": "image_to_pdf",
    }

    if ext not in conv_map:
        raise ValueError(f"Unsupported file type: {ext}")

    return convert_file(input_path, conv_map[ext], pdf_path)


# =====================================================
# 🚀 MAIN DISPATCHER
# =====================================================
def convert_file(input_path, conv_type, output_path):

    if conv_type == "merge_to_pdf":
        with tempfile.TemporaryDirectory() as temp_dir:
            pdfs = [to_pdf_if_needed(p, temp_dir) for p in input_path]
            return merge_pdfs(pdfs, output_path)

    mapping = {
        "docx_to_pdf": docx_to_pdf,
        "pptx_to_pdf": pptx_to_pdf,
        "txt_to_pdf": txt_to_pdf,
        "image_to_pdf": image_to_pdf,
        "pdf_to_docx": pdf_to_docx,
        "pdf_to_txt": pdf_to_txt,
        "pdf_to_pptx": pdf_to_pptx,
        "pdf_to_image": pdf_to_image,
    }

    if conv_type not in mapping:
        raise ValueError(f"Unsupported conversion: {conv_type}")

    return mapping[conv_type](input_path, output_path)